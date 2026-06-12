from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

from recommender import PlatformRecommender


ROOT = Path(__file__).resolve().parent
ASSET_DIR = ROOT / "ppt_assets"
OUTPUT = ROOT / "Investor_Inventor_Review_Deck.pptx"

BG = RGBColor(8, 19, 29)
PANEL = RGBColor(16, 36, 53)
PANEL_2 = RGBColor(20, 42, 62)
LINE = RGBColor(35, 64, 90)
TEXT = RGBColor(238, 247, 255)
MUTED = RGBColor(157, 181, 201)
PRIMARY = RGBColor(55, 213, 255)
SECONDARY = RGBColor(255, 177, 74)
SUCCESS = RGBColor(131, 217, 120)


def rgb(color: RGBColor) -> tuple[int, int, int]:
    return color[0], color[1], color[2]


def apply_background(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text,
    size=20,
    color=TEXT,
    bold=False,
    align=PP_ALIGN.LEFT,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.text = text
    for para in tf.paragraphs:
        para.alignment = align
        for run in para.runs:
            run.font.name = "Segoe UI"
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color
    return box


def add_panel(
    slide,
    left,
    top,
    width,
    height,
    text,
    fill=PANEL,
    line=LINE,
    size=18,
    color=TEXT,
    bold=False,
    align=PP_ALIGN.LEFT,
):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1.5)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.08)
    tf.text = text
    for para in tf.paragraphs:
        para.alignment = align
        for run in para.runs:
            run.font.name = "Segoe UI"
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color
    return shape


def add_title(slide, eyebrow, title, subtitle="") -> None:
    add_textbox(slide, Inches(0.35), Inches(0.25), Inches(6.5), Inches(0.25), eyebrow, size=14, color=PRIMARY, bold=True)
    add_textbox(slide, Inches(0.35), Inches(0.48), Inches(8.8), Inches(0.55), title, size=28, color=TEXT, bold=True)
    if subtitle:
        add_textbox(slide, Inches(0.35), Inches(1.05), Inches(9.0), Inches(0.35), subtitle, size=15, color=MUTED)


def add_full_image(slide, path: Path) -> None:
    slide.shapes.add_picture(str(path), Inches(0.10), Inches(0.15), width=Inches(13.13), height=Inches(7.10))


def add_bullets(slide, text, left, top, width, height, size=18):
    add_textbox(slide, left, top, width, height, text, size=size, color=TEXT)


recommender = PlatformRecommender()
recommender.train()
comparison = recommender.get_model_comparison()
best = comparison[0]

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def new_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide)
    return slide


# Slide 1
slide = new_slide()
add_textbox(slide, Inches(0.4), Inches(0.7), Inches(7.2), Inches(0.7), "Investor-Inventor Matchmaking System", size=28, color=TEXT, bold=True)
add_textbox(slide, Inches(0.4), Inches(1.5), Inches(7.0), Inches(0.7), "Machine learning based startup recommendation platform", size=18, color=MUTED)
add_panel(
    slide,
    Inches(0.4),
    Inches(2.35),
    Inches(6.8),
    Inches(1.2),
    "Review focus: model training, evaluation metrics, ROC/PR graphs, confusion matrices, and final model selection.",
    fill=PANEL_2,
    line=LINE,
    size=18,
)
add_panel(
    slide,
    Inches(8.0),
    Inches(1.15),
    Inches(4.1),
    Inches(2.3),
    f"Best model\n\n{best['name']}\nDecision score: {best['decision_score'] * 100:.2f}%",
    fill=PANEL,
    line=PRIMARY,
    size=20,
    bold=True,
    align=PP_ALIGN.CENTER,
)
add_textbox(slide, Inches(0.4), Inches(4.0), Inches(7.0), Inches(0.3), "Prepared for second review presentation.", size=15, color=MUTED)

# Slide 2
slide = new_slide()
add_title(slide, "Problem and Objective", "Why the system is needed")
body = (
    "• Investors need a faster way to discover relevant startup ideas.\n"
    "• Inventors need a better way to reach suitable investors.\n"
    "• Manual matching is slow, subjective, and difficult to scale.\n"
    "• The objective is to recommend strong investor-inventor pairs using machine learning.\n"
    "• The platform also supports chat, likes/dislikes, and profile review."
)
add_bullets(slide, body, Inches(0.45), Inches(1.55), Inches(7.1), Inches(3.2), size=19)
add_panel(
    slide,
    Inches(8.1),
    Inches(1.55),
    Inches(4.0),
    Inches(2.6),
    "Why it matters\n\nA good recommender reduces noise, saves time, and increases the chance of a meaningful funding connection.",
    fill=PANEL_2,
    line=LINE,
    size=18,
    bold=False,
)

# Slide 3
slide = new_slide()
add_title(slide, "Dataset and Feature Engineering", "What the model learns from")
body = (
    "• investors.csv provides funds, domain interest, location, and risk appetite.\n"
    "• inventors.csv provides startup domain, technology, location, risk, and funding need.\n"
    "• history.csv provides previous interaction outcomes.\n"
    "• Feature engineering adds domain match, location match, risk gap, affordability ratio, and text similarity.\n"
    "• Historical match statistics help the model learn which combinations worked before."
)
add_bullets(slide, body, Inches(0.45), Inches(1.55), Inches(7.0), Inches(3.5), size=18)
add_panel(
    slide,
    Inches(8.1),
    Inches(1.55),
    Inches(4.0),
    Inches(2.75),
    "Feature importance\n\nThese engineered signals are more useful than raw IDs alone because they describe how well the two profiles fit each other.",
    fill=PANEL_2,
    line=LINE,
    size=17,
)

# Slide 4
slide = new_slide()
add_full_image(slide, ASSET_DIR / "train_test_split.png")

# Slide 5
slide = new_slide()
add_title(slide, "Models Compared", "Multiple algorithms were trained on the same split")
body = (
    "• Logistic Regression: simple baseline and easy to interpret.\n"
    "• Random Forest: ensemble model that handles non-linear patterns well.\n"
    "• Gradient Boosting: sequential ensemble that often gives strong ranking performance.\n"
    "• Decision Tree: interpretable but usually less stable than ensembles.\n"
    "• Comparing all models helps us choose the most reliable recommender."
)
add_bullets(slide, body, Inches(0.45), Inches(1.55), Inches(7.0), Inches(3.4), size=18)
add_panel(
    slide,
    Inches(8.1),
    Inches(1.55),
    Inches(4.0),
    Inches(2.8),
    "Why compare models?\n\nA recommendation system should be evaluated on more than one metric. Different models can trade accuracy, precision, recall, and curve quality differently.",
    fill=PANEL_2,
    line=LINE,
    size=17,
)

# Slide 6
slide = new_slide()
add_full_image(slide, ASSET_DIR / "performance_table.png")

# Slide 7
slide = new_slide()
add_full_image(slide, ASSET_DIR / "confusion_matrices.png")

# Slide 8
slide = new_slide()
add_full_image(slide, ASSET_DIR / "roc_curve.png")

# Slide 9
slide = new_slide()
add_full_image(slide, ASSET_DIR / "pr_curve.png")

# Slide 10
slide = new_slide()
add_title(slide, "How the Best Model Is Chosen", "Model selection logic")
add_panel(
    slide,
    Inches(0.5),
    Inches(1.6),
    Inches(8.4),
    Inches(0.7),
    "Decision Score = (Accuracy + Precision + Recall + F1 + ROC AUC + PR AUC) / 6",
    fill=PANEL_2,
    line=PRIMARY,
    size=21,
    bold=True,
    align=PP_ALIGN.CENTER,
)
explain = (
    "• The best model is not chosen by accuracy alone.\n"
    "• The composite decision score gives equal importance to all six metrics.\n"
    "• ROC AUC and PR AUC are important because they show ranking quality across thresholds.\n"
    "• If two models are very close, the curve metrics break the tie.\n"
    "• In this run, Gradient Boosting edges out the others by the final decision score."
)
add_bullets(slide, explain, Inches(0.5), Inches(2.55), Inches(6.0), Inches(2.5), size=17)
ranking = (
    f"1. {comparison[0]['name']}  {comparison[0]['decision_score'] * 100:.2f}%\n"
    f"2. {comparison[1]['name']}  {comparison[1]['decision_score'] * 100:.2f}%\n"
    f"3. {comparison[2]['name']}  {comparison[2]['decision_score'] * 100:.2f}%\n"
    f"4. {comparison[3]['name']}  {comparison[3]['decision_score'] * 100:.2f}%"
)
add_panel(slide, Inches(8.0), Inches(2.55), Inches(4.2), Inches(1.8), "Final ranking\n\n" + ranking, fill=PANEL, line=LINE, size=16, bold=True)

# Slide 11
slide = new_slide()
add_title(slide, "Final Result and Demo Screenshots", "What the platform delivers")
body = (
    "• The system ranks investor-inventor pairs using machine learning.\n"
    "• Gradient Boosting is selected as the best model under the balanced decision rule.\n"
    "• The platform includes login/signup, dashboards, chat, likes/dislikes, and an equity calculator.\n"
    "• Insert a login or dashboard screenshot in the boxes below if you want a full demo slide."
)
add_bullets(slide, body, Inches(0.45), Inches(1.55), Inches(5.0), Inches(3.0), size=17)
add_panel(
    slide,
    Inches(6.0),
    Inches(1.55),
    Inches(3.1),
    Inches(1.4),
    "Insert investor dashboard screenshot here",
    fill=PANEL_2,
    line=LINE,
    size=16,
    color=MUTED,
    bold=True,
    align=PP_ALIGN.CENTER,
)
add_panel(
    slide,
    Inches(6.0),
    Inches(3.15),
    Inches(3.1),
    Inches(1.4),
    "Insert analytics dashboard screenshot here",
    fill=PANEL_2,
    line=LINE,
    size=16,
    color=MUTED,
    bold=True,
    align=PP_ALIGN.CENTER,
)
add_panel(
    slide,
    Inches(9.45),
    Inches(1.55),
    Inches(3.0),
    Inches(3.0),
    "Final takeaway\n\nThe deck shows the full ML flow: dataset, features, training, metrics, curves, and final model choice.",
    fill=PANEL,
    line=PRIMARY,
    size=17,
    bold=False,
)

if OUTPUT.exists():
    OUTPUT.unlink()
prs.save(str(OUTPUT))
print(f"PPT created: {OUTPUT}")
