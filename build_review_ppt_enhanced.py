from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

from recommender import PlatformRecommender


ROOT = Path(__file__).resolve().parent
ASSET_DIR = ROOT / "ppt_assets"
OUTPUT = ROOT / "Investor_Inventor_Review_Deck_reviewready.pptx"

USER_SCREENSHOT_1 = Path(r"C:\Users\navis\OneDrive\Pictures\Screenshots\Screenshot 2026-03-30 234135.png")
USER_SCREENSHOT_2 = Path(r"C:\Users\navis\OneDrive\Pictures\Screenshots\Screenshot 2026-03-30 234151.png")
USER_SCREENSHOT_3 = Path(r"C:\Users\navis\OneDrive\Pictures\Screenshots\Screenshot 2026-03-30 234206.png")

TRAIN_SPLIT_IMG = ASSET_DIR / "train_test_split.png"
CONFUSION_IMG = ASSET_DIR / "confusion_matrices.png"
ROC_IMG = ASSET_DIR / "roc_curve.png"
PR_IMG = ASSET_DIR / "pr_curve.png"

BG = RGBColor(8, 19, 29)
PANEL = RGBColor(16, 36, 53)
PANEL_2 = RGBColor(20, 42, 62)
LINE = RGBColor(35, 64, 90)
TEXT = RGBColor(238, 247, 255)
MUTED = RGBColor(157, 181, 201)
PRIMARY = RGBColor(55, 213, 255)
SECONDARY = RGBColor(255, 177, 74)
SUCCESS = RGBColor(131, 217, 120)
SOFT = RGBColor(42, 56, 72)


def add_background(slide, prs):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text,
    size=18,
    color=TEXT,
    bold=False,
    align=PP_ALIGN.LEFT,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.01)
    tf.margin_bottom = Inches(0.01)
    tf.text = text
    for para in tf.paragraphs:
        para.alignment = align
        for run in para.runs:
            run.font.name = "Segoe UI"
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color
    return box


def add_panel(
    slide,
    left,
    top,
    width,
    height,
    text,
    fill=PANEL,
    line=LINE,
    accent=None,
    size=17,
    color=TEXT,
    bold=False,
    align=PP_ALIGN.LEFT,
):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = accent or line
    shape.line.width = Pt(1.5)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = Inches(0.10)
    tf.margin_right = Inches(0.10)
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.08)
    tf.text = text
    for para in tf.paragraphs:
        para.alignment = align
        for run in para.runs:
            run.font.name = "Segoe UI"
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color
    return shape


def add_header(slide, eyebrow, title, subtitle=""):
    add_textbox(slide, Inches(0.35), Inches(0.22), Inches(6.5), Inches(0.22), eyebrow, size=13, color=PRIMARY, bold=True)
    add_textbox(slide, Inches(0.35), Inches(0.46), Inches(8.8), Inches(0.55), title, size=28, color=TEXT, bold=True)
    if subtitle:
        add_textbox(slide, Inches(0.35), Inches(1.00), Inches(9.6), Inches(0.30), subtitle, size=15, color=MUTED)


def image_size(path: Path):
    with Image.open(path) as img:
        return img.size


def add_picture_fit(slide, path: Path, left, top, max_width, max_height):
    iw, ih = image_size(path)
    scale = min(max_width / iw, max_height / ih)
    slide.shapes.add_picture(str(path), left, top, width=int(iw * scale), height=int(ih * scale))


def stat_card(slide, left, top, width, height, label, value, fill=PANEL, accent=SOFT):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = accent
    shape.line.width = Pt(1.2)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.text = f"{label}\n{value}"
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    tf.paragraphs[0].runs[0].font.name = "Segoe UI"
    tf.paragraphs[0].runs[0].font.size = Pt(13)
    tf.paragraphs[0].runs[0].font.color.rgb = MUTED
    tf.paragraphs[1].alignment = PP_ALIGN.LEFT
    tf.paragraphs[1].runs[0].font.name = "Segoe UI"
    tf.paragraphs[1].runs[0].font.size = Pt(21)
    tf.paragraphs[1].runs[0].font.bold = True
    tf.paragraphs[1].runs[0].font.color.rgb = TEXT
    return shape


def fit_box_text(text: str) -> str:
    return text.replace("\u2022", "-")


recommender = PlatformRecommender()
recommender.train()
comparison = recommender.get_model_comparison()
best = comparison[0]
graph = recommender.get_graph_payload()

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def new_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, prs)
    return slide


# Slide 1
slide = new_slide()
add_textbox(slide, Inches(0.40), Inches(0.65), Inches(7.2), Inches(0.6), "Investor-Inventor Matchmaking System", size=28, color=TEXT, bold=True)
add_textbox(slide, Inches(0.40), Inches(1.38), Inches(7.0), Inches(0.5), "Machine learning based startup recommendation platform", size=18, color=MUTED)
add_panel(
    slide,
    Inches(0.40),
    Inches(2.05),
    Inches(6.75),
    Inches(1.35),
    "Review focus:\nmodel training, evaluation metrics, ROC/PR graphs, confusion matrices, and final model selection.",
    fill=PANEL_2,
    line=LINE,
    size=18,
)
stat_card(slide, Inches(0.40), Inches(3.72), Inches(1.58), Inches(1.02), "Investors", str(len(recommender.investors_df)))
stat_card(slide, Inches(2.10), Inches(3.72), Inches(1.58), Inches(1.02), "Inventors", str(len(recommender.inventors_df)))
stat_card(slide, Inches(3.80), Inches(3.72), Inches(1.78), Inches(1.02), "Interactions", str(len(recommender.history_df)))
stat_card(slide, Inches(5.70), Inches(3.72), Inches(1.48), Inches(1.02), "Best model", best["name"], fill=PANEL, accent=PRIMARY)
add_panel(
    slide,
    Inches(7.95),
    Inches(1.15),
    Inches(4.10),
    Inches(2.72),
    f"Best model\n\n{best['name']}\nDecision score: {best['decision_score'] * 100:.2f}%\nROC AUC: {best['roc_auc'] * 100:.2f}%\nPR AUC: {best['pr_auc'] * 100:.2f}%",
    fill=PANEL,
    line=PRIMARY,
    size=18,
    bold=True,
    align=PP_ALIGN.CENTER,
)
add_textbox(slide, Inches(0.40), Inches(5.10), Inches(7.0), Inches(0.25), "Prepared for second review presentation.", size=15, color=MUTED)


# Slide 2
slide = new_slide()
add_header(slide, "Problem and Objective", "Why the system is needed")
add_panel(
    slide,
    Inches(0.45),
    Inches(1.52),
    Inches(4.00),
    Inches(4.55),
    "Problem statement\n\n- Investors need a faster way to discover relevant startup ideas.\n- Inventors need a better way to reach suitable investors.\n- Manual matching is slow, subjective, and difficult to scale.\n- A common platform can reduce noise and improve discovery.",
    fill=PANEL_2,
    size=18,
)
add_panel(
    slide,
    Inches(4.65),
    Inches(1.52),
    Inches(3.55),
    Inches(4.55),
    "Project objective\n\n- Recommend strong investor-inventor pairs using machine learning.\n- Show the best matches with scores and explanations.\n- Let users like, dislike, and chat through the platform.\n- Support a real Shark Tank style workflow.",
    fill=PANEL,
    accent=SECONDARY,
    size=17,
)
add_panel(
    slide,
    Inches(8.35),
    Inches(1.52),
    Inches(4.50),
    Inches(4.55),
    "Workflow\n\n1. User logs in or signs up.\n2. Investor sees recommended ideas.\n3. Inventor can review interested investors.\n4. Both sides can open the shared chat.\n5. Investor can shortlist or reject ideas.\n6. Inventor can share achievements and patents.",
    fill=PANEL_2,
    accent=SUCCESS,
    size=16,
)
add_panel(
    slide,
    Inches(0.45),
    Inches(6.26),
    Inches(12.35),
    Inches(0.55),
    "Why it matters: a good recommender reduces search time, improves match quality, and gives a more structured funding conversation.",
    fill=PANEL,
    line=LINE,
    size=15,
    align=PP_ALIGN.CENTER,
)


# Slide 3
slide = new_slide()
add_header(slide, "Dataset and Feature Engineering", "What the model learns from")
add_panel(
    slide,
    Inches(0.45),
    Inches(1.52),
    Inches(4.00),
    Inches(4.65),
    "Data sources\n\n- investors.csv: funds, domain interest, location, risk appetite.\n- inventors.csv: startup domain, technology, funding need, location, risk level.\n- history.csv: past interaction outcomes.\n- Augmented rows were added only to strengthen training, not testing.",
    fill=PANEL_2,
    size=17,
)
add_panel(
    slide,
    Inches(4.65),
    Inches(1.52),
    Inches(4.10),
    Inches(4.65),
    "Engineered features\n\n- domain match\n- location match\n- risk gap\n- affordability ratio\n- text similarity\n- investor/idea historical rates",
    fill=PANEL,
    accent=PRIMARY,
    size=18,
)
add_panel(
    slide,
    Inches(8.85),
    Inches(1.52),
    Inches(3.95),
    Inches(4.65),
    "Why these features help\n\nThe model does not rely on raw IDs alone. It learns fit signals that are closer to how a human would judge a good funding match: same domain, acceptable risk, viable funding, and meaningful interest overlap.",
    fill=PANEL_2,
    accent=SECONDARY,
    size=16,
)
add_panel(
    slide,
    Inches(0.45),
    Inches(6.26),
    Inches(12.35),
    Inches(0.55),
    "Core idea: turn two profile tables into one pair-wise feature set, then learn which combinations historically produced strong interactions.",
    fill=PANEL,
    line=LINE,
    size=15,
    align=PP_ALIGN.CENTER,
)


# Slide 4
slide = new_slide()
add_header(slide, "Training and Validation", "Train/test split diagram")
add_picture_fit(slide, TRAIN_SPLIT_IMG, Inches(0.30), Inches(1.55), Inches(8.20), Inches(5.55))
add_panel(
    slide,
    Inches(8.75),
    Inches(1.55),
    Inches(4.25),
    Inches(1.55),
    "75% training / 25% testing\nThe training set teaches the model. The test set stays unseen until final evaluation.",
    fill=PANEL_2,
    size=17,
)
add_panel(
    slide,
    Inches(8.75),
    Inches(3.22),
    Inches(4.25),
    Inches(1.55),
    "Fair comparison\nThe same split is used for every model so the comparison is consistent and unbiased.",
    fill=PANEL,
    accent=SECONDARY,
    size=17,
)
add_panel(
    slide,
    Inches(8.75),
    Inches(4.89),
    Inches(4.25),
    Inches(1.55),
    "No test leakage\nSynthetic rows help training only. They do not decide the reported test score.",
    fill=PANEL_2,
    accent=SUCCESS,
    size=17,
)


# Slide 5
slide = new_slide()
add_header(slide, "Model Comparison Evidence", "Performance table from the dashboard")
add_picture_fit(slide, USER_SCREENSHOT_1, Inches(0.28), Inches(1.55), Inches(6.95), Inches(5.55))
add_panel(
    slide,
    Inches(7.45),
    Inches(1.55),
    Inches(5.45),
    Inches(1.50),
    "How to read this table\n\nEach row is one model. The columns show Accuracy, Precision, Recall, F1, ROC AUC, PR AUC, and the final decision score.",
    fill=PANEL_2,
    size=16,
)
add_panel(
    slide,
    Inches(7.45),
    Inches(3.17),
    Inches(5.45),
    Inches(1.55),
    "Key observation\n\nGradient Boosting has the strongest ROC AUC and PR AUC. Logistic Regression has slightly higher accuracy, but accuracy alone is not enough for ranking recommendations.",
    fill=PANEL,
    accent=PRIMARY,
    size=16,
)
add_panel(
    slide,
    Inches(7.45),
    Inches(4.84),
    Inches(5.45),
    Inches(1.55),
    f"Review line\n\nBest model: {best['name']}\nDecision score: {best['decision_score'] * 100:.2f}%\nThe gap is tiny, so the winner is chosen by the full composite metric, not by one number only.",
    fill=PANEL_2,
    accent=SECONDARY,
    size=15,
)


# Slide 6
slide = new_slide()
add_header(slide, "Performance Graphs and Demand Insight", "Metric bars plus catalog insight")
add_picture_fit(slide, USER_SCREENSHOT_2, Inches(0.28), Inches(1.55), Inches(7.30), Inches(5.55))
add_panel(
    slide,
    Inches(7.80),
    Inches(1.55),
    Inches(5.10),
    Inches(1.50),
    "Metric bars\nThis section summarizes each model in one line so the review panel can quickly compare accuracy, precision, recall, F1, ROC AUC, PR AUC, and decision score.",
    fill=PANEL_2,
    size=15,
)
add_panel(
    slide,
    Inches(7.80),
    Inches(3.15),
    Inches(5.10),
    Inches(1.55),
    "Catalog insight\nThe right side shows which domains and technologies are most common in the dataset. This helps explain where the data is dense and what sectors are prominent.",
    fill=PANEL,
    accent=SUCCESS,
    size=15,
)
add_panel(
    slide,
    Inches(7.80),
    Inches(4.82),
    Inches(5.10),
    Inches(1.58),
    "Why this slide matters\nIt connects model performance with the underlying dataset distribution, which makes the platform feel more complete than a pure ML demo.",
    fill=PANEL_2,
    accent=PRIMARY,
    size=15,
)


# Slide 7
slide = new_slide()
add_header(slide, "Comparison Matrix", "Metric heatmap")
add_picture_fit(slide, USER_SCREENSHOT_3, Inches(0.26), Inches(1.18), Inches(8.25), Inches(6.05))
add_panel(
    slide,
    Inches(8.75),
    Inches(1.18),
    Inches(4.18),
    Inches(2.10),
    "How to compare the models\n\n- Accuracy shows overall correctness.\n- Precision shows how many predicted matches are truly good.\n- Recall shows how many real good matches were found.\n- F1 balances precision and recall.",
    fill=PANEL_2,
    size=16,
)
add_panel(
    slide,
    Inches(8.75),
    Inches(3.44),
    Inches(4.18),
    Inches(2.25),
    "What the matrix tells us\n\nGradient Boosting and Logistic Regression are very close. Random Forest is also strong. Decision Tree is weaker. The final choice is based on the full score, not a single metric.",
    fill=PANEL,
    accent=SECONDARY,
    size=16,
)
add_panel(
    slide,
    Inches(8.75),
    Inches(5.85),
    Inches(4.18),
    Inches(1.00),
    f"Decision score leader: {best['name']} ({best['decision_score'] * 100:.2f}%)",
    fill=PANEL_2,
    accent=PRIMARY,
    size=16,
    align=PP_ALIGN.CENTER,
)


# Slide 8
slide = new_slide()
add_header(slide, "Confusion Matrices", "Correct vs incorrect predictions")
add_picture_fit(slide, CONFUSION_IMG, Inches(0.30), Inches(1.50), Inches(8.15), Inches(5.65))
add_panel(
    slide,
    Inches(8.75),
    Inches(1.50),
    Inches(4.20),
    Inches(1.55),
    "TP / TN / FP / FN\nTP = correct good match\nTN = correct bad match\nFP = wrong positive recommendation\nFN = missed good opportunity",
    fill=PANEL_2,
    size=15,
)
add_panel(
    slide,
    Inches(8.75),
    Inches(3.20),
    Inches(4.20),
    Inches(1.55),
    "How to conclude it\nThe top models have similar confusion counts. That is why the ROC and PR curves are important for the final decision.",
    fill=PANEL,
    accent=SUCCESS,
    size=16,
)
add_panel(
    slide,
    Inches(8.75),
    Inches(4.90),
    Inches(4.20),
    Inches(1.55),
    "Interpretation\nDecision Tree makes more mistakes. The ensemble models keep the mistakes lower and the ranking quality higher.",
    fill=PANEL_2,
    accent=SECONDARY,
    size=15,
)


# Slide 9
slide = new_slide()
add_header(slide, "ROC Curve", "True positive rate vs false positive rate")
add_picture_fit(slide, ROC_IMG, Inches(0.25), Inches(1.50), Inches(8.10), Inches(5.65))
add_panel(
    slide,
    Inches(8.70),
    Inches(1.50),
    Inches(4.25),
    Inches(1.55),
    "How to read ROC\nX-axis = False Positive Rate\nY-axis = True Positive Rate\nCurves closer to the top-left are stronger.",
    fill=PANEL_2,
    size=15,
)
add_panel(
    slide,
    Inches(8.70),
    Inches(3.20),
    Inches(4.25),
    Inches(1.60),
    f"What we see here\nGradient Boosting leads with ROC AUC {comparison[0]['roc_auc'] * 100:.2f}.\nLogistic Regression is close, and Random Forest follows behind.",
    fill=PANEL,
    accent=PRIMARY,
    size=15,
)
add_panel(
    slide,
    Inches(8.70),
    Inches(4.95),
    Inches(4.25),
    Inches(1.45),
    "Why it matters\nThe ROC curve shows the model's ability to separate good matches from bad matches across thresholds.",
    fill=PANEL_2,
    accent=SECONDARY,
    size=15,
)


# Slide 10
slide = new_slide()
add_header(slide, "Precision-Recall Curve", "Precision vs recall across the top models")
add_picture_fit(slide, PR_IMG, Inches(0.25), Inches(1.50), Inches(8.10), Inches(5.65))
add_panel(
    slide,
    Inches(8.70),
    Inches(1.50),
    Inches(4.25),
    Inches(1.55),
    "How to read PR\nX-axis = Recall\nY-axis = Precision\nHigher curves are better when we care about top match quality.",
    fill=PANEL_2,
    size=15,
)
add_panel(
    slide,
    Inches(8.70),
    Inches(3.20),
    Inches(4.25),
    Inches(1.60),
    f"What we see here\nGradient Boosting leads with PR AUC {comparison[0]['pr_auc'] * 100:.2f}.\nThis is the strongest curve-based signal in the deck.",
    fill=PANEL,
    accent=SUCCESS,
    size=15,
)
add_panel(
    slide,
    Inches(8.70),
    Inches(4.95),
    Inches(4.25),
    Inches(1.45),
    f"Baseline note\nThe positive-class rate in the test split is {graph['positive_rate']:.2f}%. Curves above that baseline are better than a naive positive guess.",
    fill=PANEL_2,
    accent=SECONDARY,
    size=15,
)


# Slide 11
slide = new_slide()
add_header(slide, "Model Selection Logic", "How the winner is decided")
add_panel(
    slide,
    Inches(0.45),
    Inches(1.35),
    Inches(8.25),
    Inches(0.72),
    "Decision score = (Accuracy + Precision + Recall + F1 + ROC AUC + PR AUC) / 6",
    fill=PANEL_2,
    line=PRIMARY,
    size=21,
    bold=True,
    align=PP_ALIGN.CENTER,
)
add_panel(
    slide,
    Inches(0.45),
    Inches(2.25),
    Inches(4.00),
    Inches(3.35),
    "Why this rule is used\n\n- Accuracy alone can hide ranking quality.\n- Precision and recall show the quality of positive matches.\n- ROC AUC and PR AUC show threshold-based discrimination.\n- Equal weighting keeps the model choice balanced and explainable.",
    fill=PANEL_2,
    size=16,
)
rank_text = (
    f"Final ranking\n\n"
    f"1. {comparison[0]['name']}  {comparison[0]['decision_score'] * 100:.2f}%\n"
    f"2. {comparison[1]['name']}  {comparison[1]['decision_score'] * 100:.2f}%\n"
    f"3. {comparison[2]['name']}  {comparison[2]['decision_score'] * 100:.2f}%\n"
    f"4. {comparison[3]['name']}  {comparison[3]['decision_score'] * 100:.2f}%"
)
add_panel(
    slide,
    Inches(4.65),
    Inches(2.25),
    Inches(4.05),
    Inches(3.35),
    rank_text,
    fill=PANEL,
    accent=SECONDARY,
    size=16,
    bold=True,
)
add_panel(
    slide,
    Inches(8.95),
    Inches(2.25),
    Inches(3.95),
    Inches(3.35),
    f"Why Gradient Boosting wins\n\n- It has the strongest ROC AUC and PR AUC.\n- Its class metrics remain very strong.\n- The margin over Logistic Regression is tiny, but real.\n- For recommendation ranking, curve quality is a key advantage.",
    fill=PANEL_2,
    accent=SUCCESS,
    size=16,
)
add_panel(
    slide,
    Inches(0.45),
    Inches(5.92),
    Inches(12.45),
    Inches(0.68),
    "If asked in review: the models are close, but Gradient Boosting is highlighted because it offers the best overall balance for a recommendation system, not just the highest single metric.",
    fill=PANEL,
    line=LINE,
    size=15,
    align=PP_ALIGN.CENTER,
)


# Slide 12
slide = new_slide()
add_header(slide, "Platform Features", "What the full application does")
add_panel(
    slide,
    Inches(0.45),
    Inches(1.55),
    Inches(4.05),
    Inches(3.95),
    "Investor side\n\n- Login and signup\n- View recommended ideas\n- Like / dislike matches\n- Open shared chat\n- Search inventors from the dataset\n- Use equity calculator",
    fill=PANEL_2,
    size=17,
)
add_panel(
    slide,
    Inches(4.80),
    Inches(1.55),
    Inches(4.05),
    Inches(3.95),
    "Inventor side\n\n- Add profile details\n- Add achievements and patents\n- View interested investors\n- Respond in shared chat\n- Track feedback from investors\n- Improve visibility through the platform",
    fill=PANEL,
    accent=PRIMARY,
    size=17,
)
add_panel(
    slide,
    Inches(9.15),
    Inches(1.55),
    Inches(3.75),
    Inches(3.95),
    "Why this is more than ML\n\nThe platform is a middleware layer between investors and inventors. It combines recommendation, communication, and profile sharing in one workflow.",
    fill=PANEL_2,
    accent=SECONDARY,
    size=16,
)
add_panel(
    slide,
    Inches(0.45),
    Inches(5.80),
    Inches(12.45),
    Inches(0.72),
    "User journey: login -> recommendations -> feedback -> chat -> shortlist -> possible funding discussion.",
    fill=PANEL,
    line=LINE,
    size=15,
    align=PP_ALIGN.CENTER,
)


# Slide 13
slide = new_slide()
add_header(slide, "Conclusion and Future Scope", "What faculty should remember")
add_panel(
    slide,
    Inches(0.45),
    Inches(1.52),
    Inches(4.05),
    Inches(4.55),
    "Key takeaways\n\n- The platform matches investors and inventors using machine learning.\n- We compared multiple models on the same split.\n- We used both class metrics and curve metrics.\n- Gradient Boosting is selected as the best balanced model.\n- The app includes chat and investor/inventor workflows.",
    fill=PANEL_2,
    size=17,
)
add_panel(
    slide,
    Inches(4.80),
    Inches(1.52),
    Inches(4.05),
    Inches(4.55),
    f"Final model summary\n\n{best['name']}\n\nAccuracy: {best['accuracy'] * 100:.2f}%\nPrecision: {best['precision'] * 100:.2f}%\nRecall: {best['recall'] * 100:.2f}%\nF1: {best['f1'] * 100:.2f}%\nROC AUC: {best['roc_auc'] * 100:.2f}%\nPR AUC: {best['pr_auc'] * 100:.2f}%",
    fill=PANEL,
    accent=PRIMARY,
    size=16,
)
add_panel(
    slide,
    Inches(9.15),
    Inches(1.52),
    Inches(3.75),
    Inches(4.55),
    "Future scope\n\n- Add real user data\n- Improve chat notifications\n- Add richer explainability\n- Track investor feedback over time\n- Add meeting/request workflow",
    fill=PANEL_2,
    accent=SECONDARY,
    size=16,
)
add_panel(
    slide,
    Inches(0.45),
    Inches(6.25),
    Inches(12.45),
    Inches(0.52),
    "Thank you - the system is ready for demonstration and review.",
    fill=PANEL,
    line=LINE,
    size=15,
    align=PP_ALIGN.CENTER,
)


if OUTPUT.exists():
    OUTPUT.unlink()
prs.save(str(OUTPUT))
print(f"PPT created: {OUTPUT}")
